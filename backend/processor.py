import re
from openai_utils import generate_slide_script
OUTPUTS_CONTAINER = "outputs"

# Global jobs dictionary for tracking progress (if used in Flask context)
jobs = {}

# Clean text utility
def clean_text(text):
    if not text:
        return ""
    # Remove control/invisible characters
    text = re.sub(r'[\u0000-\u001F\u007F-\u009F\u200b\u000b]', '', text)
    # Replace multiple newlines with a single newline
    text = re.sub(r'\n+', '\n', text)
    return text.strip()
import os, subprocess, tempfile, uuid, json, shutil

# Ensure subprocess calls don't create windows on Windows
if os.name == 'nt':
    # Set default creation flags for all subprocess calls in this module
    subprocess._USE_POSIX_SPAWN = False

from datetime import datetime, timedelta
from azure.storage.blob import generate_blob_sas, BlobSasPermissions
from pptx import Presentation
from blob_utils import download_to_path, upload_file_path, upload_file_bytes, SLIDES_CONTAINER

def convert_pptx_to_pngs(local_pptx_path, out_dir):
    # Check file size before processing
    file_size = os.path.getsize(local_pptx_path)
    file_size_mb = file_size / (1024 * 1024)
    print(f"[convert_pptx_to_pngs] Processing file size: {file_size_mb:.1f} MB")
    
    if file_size_mb > 50:
        print(f"[convert_pptx_to_pngs] WARNING: Large file ({file_size_mb:.1f} MB) may take longer to process")
    
    # Validate the PowerPoint file first
    try:
        print(f"[convert_pptx_to_pngs] Validating PowerPoint file...")
        if not os.path.exists(local_pptx_path):
            raise Exception(f"Input file does not exist: {local_pptx_path}")
        
        if not os.path.isfile(local_pptx_path):
            raise Exception(f"Input path is not a file: {local_pptx_path}")
            
        # Check if file is accessible
        with open(local_pptx_path, 'rb') as f:
            f.read(10)  # Try to read first 10 bytes
            
        ppt = Presentation(local_pptx_path)
        slide_count = len(ppt.slides)
        print(f"[convert_pptx_to_pngs] File validation successful: {slide_count} slides found")
        ppt = None  # Release the file handle
    except Exception as e:
        print(f"[convert_pptx_to_pngs] PowerPoint file validation failed: {e}")
        raise Exception(f"Invalid or corrupted PowerPoint file: {e}")
    
    # Validate output directory
    if not os.path.exists(out_dir):
        try:
            os.makedirs(out_dir, exist_ok=True)
        except Exception as e:
            print(f"[convert_pptx_to_pngs] Failed to create output directory: {e}")
            raise Exception(f"Cannot create output directory: {e}")
    
    # Test write permissions
    test_file = os.path.join(out_dir, 'test_write.tmp')
    try:
        with open(test_file, 'w') as f:
            f.write('test')
        os.remove(test_file)
        print(f"[convert_pptx_to_pngs] Output directory write test successful")
    except Exception as e:
        print(f"[convert_pptx_to_pngs] Output directory write test failed: {e}")
        raise Exception(f"Cannot write to output directory: {e}")
    
    # Convert PPTX -> PDF using LibreOffice
    print(f"[convert_pptx_to_pngs] Looking for LibreOffice executable...")
    exe = shutil.which('soffice') or shutil.which('libreoffice')
    if not exe:
        possible_path = r'C:\Program Files\LibreOffice\program\soffice.exe'
        if os.path.exists(possible_path):
            exe = possible_path
        else:
            print("[convert_pptx_to_pngs] LibreOffice executable not found.")
            raise FileNotFoundError("LibreOffice executable not found.")
    print(f"[convert_pptx_to_pngs] Using LibreOffice at: {exe}")
    
    # Convert paths to short format on Windows to avoid path length issues
    if os.name == 'nt':
        try:
            import ctypes
            from ctypes import wintypes
            _GetShortPathNameW = ctypes.windll.kernel32.GetShortPathNameW
            _GetShortPathNameW.argtypes = [wintypes.LPCWSTR, wintypes.LPWSTR, wintypes.DWORD]
            _GetShortPathNameW.restype = wintypes.DWORD
            
            def get_short_path(long_path):
                output_buf_size = 260
                output_buf = ctypes.create_unicode_buffer(output_buf_size)
                ret = _GetShortPathNameW(long_path, output_buf, output_buf_size)
                if ret:
                    return output_buf.value
                return long_path
            
            short_input_path = get_short_path(local_pptx_path)
            short_output_dir = get_short_path(out_dir)
            print(f"[convert_pptx_to_pngs] Using short paths - Input: {short_input_path}, Output: {short_output_dir}")
            local_pptx_path = short_input_path
            out_dir = short_output_dir
        except Exception as e:
            print(f"[convert_pptx_to_pngs] Could not get short paths: {e}, using original paths")
    
    # Kill any existing LibreOffice processes that might be stuck
    try:
        if os.name == 'nt':  # Windows
            subprocess.run(['taskkill', '/f', '/im', 'soffice.exe'], capture_output=True, creationflags=subprocess.CREATE_NO_WINDOW)
            subprocess.run(['taskkill', '/f', '/im', 'soffice.bin'], capture_output=True, creationflags=subprocess.CREATE_NO_WINDOW)
    except:
        pass  # Ignore errors if no processes to kill
    
    # Enhanced LibreOffice command for better PDF quality
    cmd = [
        exe, 
        '--headless',
        '--invisible', 
        '--nodefault',
        '--nologo',
        '--nolockcheck',
        '--convert-to', 'pdf',
        local_pptx_path,
        '--outdir', out_dir
    ]
    
    print(f"[convert_pptx_to_pngs] LibreOffice executable: {exe}")
    print(f"[convert_pptx_to_pngs] Input file: {local_pptx_path}")
    print(f"[convert_pptx_to_pngs] Input file exists: {os.path.exists(local_pptx_path)}")
    print(f"[convert_pptx_to_pngs] Output directory: {out_dir}")
    print(f"[convert_pptx_to_pngs] Output directory exists: {os.path.exists(out_dir)}")
    print(f"[convert_pptx_to_pngs] Running command: {' '.join(cmd)}")
    
    try:
        # First attempt with enhanced options
        process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, creationflags=subprocess.CREATE_NO_WINDOW)
        stdout, stderr = process.communicate(timeout=300)
        
        if process.returncode != 0:
            stdout_text = stdout.decode('utf-8', errors='ignore') if stdout else 'None'
            stderr_text = stderr.decode('utf-8', errors='ignore') if stderr else 'None'
            print(f"[convert_pptx_to_pngs] First attempt failed with return code {process.returncode}")
            print(f"[convert_pptx_to_pngs] stdout: {stdout_text}")
            print(f"[convert_pptx_to_pngs] stderr: {stderr_text}")
            
            # Try a simpler command as fallback
            print(f"[convert_pptx_to_pngs] Trying simpler LibreOffice command...")
            simple_cmd = [exe, '--headless', '--convert-to', 'pdf', local_pptx_path, '--outdir', out_dir]
            print(f"[convert_pptx_to_pngs] Running simple command: {' '.join(simple_cmd)}")
            
            process = subprocess.Popen(simple_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, creationflags=subprocess.CREATE_NO_WINDOW)
            stdout, stderr = process.communicate(timeout=300)
            
            if process.returncode != 0:
                stdout_text = stdout.decode('utf-8', errors='ignore') if stdout else 'None'
                stderr_text = stderr.decode('utf-8', errors='ignore') if stderr else 'None'
                print(f"[convert_pptx_to_pngs] Simple command also failed with return code {process.returncode}")
                print(f"[convert_pptx_to_pngs] stdout: {stdout_text}")
                print(f"[convert_pptx_to_pngs] stderr: {stderr_text}")
                print(f"[convert_pptx_to_pngs] Input file exists: {os.path.exists(local_pptx_path)}")
                print(f"[convert_pptx_to_pngs] Input file size: {os.path.getsize(local_pptx_path) if os.path.exists(local_pptx_path) else 'N/A'}")
                print(f"[convert_pptx_to_pngs] Output directory exists: {os.path.exists(out_dir)}")
                print(f"[convert_pptx_to_pngs] Output directory contents: {os.listdir(out_dir) if os.path.exists(out_dir) else 'N/A'}")
                
                # Common LibreOffice error solutions
                if "Permission denied" in stderr_text or "Access denied" in stderr_text:
                    raise Exception(f"LibreOffice permission denied - file may be locked or in use. Return code: {process.returncode}")
                elif "java.lang.OutOfMemoryError" in stderr_text:
                    raise Exception(f"LibreOffice out of memory - file too large or complex. Return code: {process.returncode}")
                elif "format" in stderr_text.lower() or "corrupt" in stderr_text.lower():
                    raise Exception(f"LibreOffice cannot read file - may be corrupted or unsupported format. Return code: {process.returncode}")
                else:
                    raise Exception(f"LibreOffice conversion failed with return code {process.returncode}. stdout: {stdout_text[:200]}... stderr: {stderr_text[:200]}...")
            else:
                print(f"[convert_pptx_to_pngs] Simple command succeeded!")
                
    except subprocess.TimeoutExpired:
        print("[convert_pptx_to_pngs] LibreOffice conversion timed out after 300 seconds!")
        # Kill the process
        process.kill()
        process.wait()
        raise Exception("LibreOffice conversion timed out after 5 minutes")
    except Exception as e:
        print(f"[convert_pptx_to_pngs] LibreOffice conversion failed: {e}")
        raise

    base = os.path.splitext(os.path.basename(local_pptx_path))[0]
    pdf_path = os.path.join(out_dir, base + '.pdf')
    if not os.path.exists(pdf_path):
        print("[convert_pptx_to_pngs] PDF conversion failed: PDF not found.")
        raise Exception("PDF conversion failed.")

    print(f"[convert_pptx_to_pngs] Looking for pdftoppm executable...")
    out_prefix = os.path.join(out_dir, "slide")
    pdftoppm_exe = shutil.which('pdftoppm')
    if not pdftoppm_exe:
        possible_poppler = r'C:\Program Files\poppler-24.02.0\Library\bin\pdftoppm.exe'
        if os.path.exists(possible_poppler):
            pdftoppm_exe = possible_poppler
        else:
            print("[convert_pptx_to_pngs] pdftoppm executable not found.")
            raise FileNotFoundError("pdftoppm executable not found.")
    print(f"[convert_pptx_to_pngs] Using pdftoppm at: {pdftoppm_exe}")
    try:
        subprocess.check_call([pdftoppm_exe, '-png', pdf_path, out_prefix], timeout=600, creationflags=subprocess.CREATE_NO_WINDOW)
    except subprocess.TimeoutExpired:
        print("[convert_pptx_to_pngs] pdftoppm conversion timed out!")
        raise
    except Exception as e:
        print(f"[convert_pptx_to_pngs] pdftoppm conversion failed: {e}")
        raise

    # Collect slide PNGs
    files = sorted([os.path.join(out_dir, f) for f in os.listdir(out_dir) if f.startswith('slide') and f.endswith('.png')])
    print(f"[convert_pptx_to_pngs] Found {len(files)} PNG slides.")
    return files

def extract_notes(local_pptx_path):
    prs = Presentation(local_pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title else f"Slide {i}"
        notes = ""
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"index": i, "title": title, "notes": notes})
    return slides

def process_presentation_from_blob(blob_container: str, blob_name: str, avatar: str = "harry"):
    # Avatar presets
    AVATAR_PRESETS = {
        "harry": {
            "character": "Harry",
            "style": "business",
            "voice": "en-IN-ArjunNeural"
        },
        "lisa": {
            "character": "lisa",
            "style": "casual-sitting",
            "voice": "en-US-AvaMultilingualNeural"
        }
    }
    avatar_cfg = AVATAR_PRESETS.get(avatar, AVATAR_PRESETS["harry"])
    tmpdir = tempfile.mkdtemp()
    local_ppt = os.path.join(tmpdir, blob_name)
    download_to_path(blob_container, blob_name, local_ppt)


    pngs = convert_pptx_to_pngs(local_ppt, tmpdir)
    slide_notes = extract_notes(local_ppt)

    presentation_id = str(uuid.uuid4())
    slides_info = []
    narrations = []

    # --- Progress/ETA logic ---
    import time
    from flask import has_request_context, g
    # Try to get job_id from context if available
    job_id = getattr(g, 'job_id', None) if has_request_context() else None
    total_slides = len(pngs)
    start_time = time.time()

    for i, png_path in enumerate(pngs, start=1):
        blob_name_out = f"{presentation_id}/slide-{i}.png"
        url = upload_file_path(SLIDES_CONTAINER, blob_name_out, png_path)
        notes = clean_text(slide_notes[i-1]['notes'] if i-1 < len(slide_notes) else "")
        title = clean_text(slide_notes[i-1]['title'] if i-1 < len(slide_notes) else f"Slide {i}")

        narration = ""
        narration_blob_name = f"{presentation_id}/slide-{i}-narration.txt"
        print(f"[Slide {i}] Title: {title}")
        print(f"[Slide {i}] Notes: {notes}")

        try:
            narration_input = notes
            if not notes:
                narration_input = f"This slide is titled '{title}'. Please provide a brief spoken summary based on the title."
            narration = generate_slide_script(title, narration_input, avatar_cfg.get("voice")) if "voice" in avatar_cfg else generate_slide_script(title, narration_input)
            narration = re.sub(r'[\u0000-\u001F\u007F-\u009F\u200b\u000b\u00a0]', '', narration)
            print(f"[Slide {i}] Narration: {narration}")
            upload_file_bytes(OUTPUTS_CONTAINER, narration_blob_name, narration.encode('utf-8'))
        except Exception as e:
            print(f"[Slide {i}] Narration generation/upload failed: {e}")
            narration_blob_name = ""
            narration = ""

        narrations.append({
            "index": i,
            "title": title,
            "narration": narration,
            "image_blob": url,
            "image_blob_name": blob_name_out
        })

        slides_info.append({
            "index": i,
            "image_blob": url,
            "image_blob_name": blob_name_out,
            "title": title,
            "notes": notes,
            "narration_blob": narration_blob_name if narration else ""
        })

        # --- Update progress/ETA in jobs dict if job_id is available ---
        if 'jobs' in globals() and hasattr(g, 'job_id'):
            elapsed = time.time() - start_time
            slides_done = i
            percent = int((slides_done / total_slides) * 100)
            avg_time_per_slide = elapsed / slides_done if slides_done > 0 else 0
            eta = int(avg_time_per_slide * (total_slides - slides_done))
            jobs[g.job_id]["progress"] = percent
            jobs[g.job_id]["eta"] = eta
            jobs[g.job_id]["slides_done"] = slides_done
            jobs[g.job_id]["total_slides"] = total_slides


    metadata = {"presentation_id": presentation_id, "slides": slides_info}

    # Upload metadata JSON to Blob
    meta_blob = f"{presentation_id}/metadata.json"
    upload_file_bytes(SLIDES_CONTAINER, meta_blob, json.dumps(metadata).encode('utf-8'))

    # Upload narration JSON to outputs container
    narration_json_blob = f"{presentation_id}/narration.json"
    # Always include both narrations and avatar config
    narration_json_data = {
        "narrations": narrations,
        "avatar": avatar_cfg
    }
    upload_file_bytes(OUTPUTS_CONTAINER, narration_json_blob, json.dumps(narration_json_data, ensure_ascii=False).encode('utf-8'))

    # Example: If you have downstream avatar video synthesis, use avatar_cfg["character"] and avatar_cfg["style"]
    # e.g., submit_synthesis(job_id, avatar_cfg["character"], avatar_cfg["style"], avatar_cfg["voice"])


    # With Azure AD, SAS token generation using account_key is not supported. Just upload narration JSON and return metadata.
    return metadata
